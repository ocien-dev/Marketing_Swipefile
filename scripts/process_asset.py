#!/usr/bin/env python
"""Process registered complementary assets into normalized content segments."""

from __future__ import annotations

import argparse
import csv
import html
import json
import re
from pathlib import Path
from typing import Any

from youtube_common import write_json


def load_json(path: Path) -> dict[str, Any]:
    with path.open("r", encoding="utf-8") as file:
        return json.load(file)


def clean_text(text: str) -> str:
    text = html.unescape(text)
    text = re.sub(r"<[^>]+>", " ", text)
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def split_text_blocks(text: str) -> list[tuple[str | None, str]]:
    lines = text.splitlines()
    blocks: list[tuple[str | None, str]] = []
    current_title: str | None = None
    current_lines: list[str] = []

    def flush() -> None:
        nonlocal current_lines
        joined = clean_text("\n".join(current_lines))
        if joined:
            blocks.append((current_title, joined))
        current_lines = []

    for line in lines:
        stripped = line.strip()
        if not stripped:
            flush()
            continue
        if stripped.startswith("#"):
            flush()
            current_title = stripped.lstrip("#").strip() or None
            continue
        current_lines.append(stripped)
    flush()

    if not blocks:
        cleaned = clean_text(text)
        if cleaned:
            blocks.append((None, cleaned))
    return blocks


def text_segments(metadata: dict[str, Any], file_path: Path) -> list[dict[str, Any]]:
    text = file_path.read_text(encoding="utf-8", errors="replace")
    blocks = split_text_blocks(text)
    asset_id = metadata["asset_id"]
    segments = []
    for index, (section_title, block_text) in enumerate(blocks):
        segments.append(
            {
                "segment_id": f"{asset_id}-asset-{index + 1:04d}",
                "segment_index": index,
                "source_kind": "asset",
                "start_seconds": None,
                "end_seconds": None,
                "page_number": None,
                "sheet_name": None,
                "cell_range": None,
                "slide_number": None,
                "section_title": section_title,
                "text_original": block_text,
                "text_ptbr": None,
                "language": metadata.get("language_original"),
            }
        )
    return segments


def csv_segments(metadata: dict[str, Any], file_path: Path) -> list[dict[str, Any]]:
    asset_id = metadata["asset_id"]
    segments = []
    with file_path.open("r", encoding="utf-8-sig", newline="") as file:
        reader = csv.DictReader(file)
        headers = reader.fieldnames or []
        for row_index, row in enumerate(reader, start=2):
            parts = [f"{header}: {row.get(header, '')}" for header in headers]
            text = clean_text("; ".join(parts))
            if not text:
                continue
            segments.append(
                {
                    "segment_id": f"{asset_id}-asset-{len(segments) + 1:04d}",
                    "segment_index": len(segments),
                    "source_kind": "asset",
                    "start_seconds": None,
                    "end_seconds": None,
                    "page_number": None,
                    "sheet_name": file_path.stem,
                    "cell_range": f"A{row_index}:{chr(ord('A') + max(len(headers) - 1, 0))}{row_index}" if headers else None,
                    "slide_number": None,
                    "section_title": f"CSV row {row_index}",
                    "text_original": text,
                    "text_ptbr": None,
                    "language": metadata.get("language_original"),
                }
            )
    return segments


def pdf_segments(metadata: dict[str, Any], file_path: Path) -> list[dict[str, Any]]:
    import pdfplumber

    asset_id = metadata["asset_id"]
    segments = []
    with pdfplumber.open(str(file_path)) as pdf:
        for page_index, page in enumerate(pdf.pages, start=1):
            text = clean_text(page.extract_text() or "")
            if not text:
                continue
            segments.append(
                {
                    "segment_id": f"{asset_id}-asset-{len(segments) + 1:04d}",
                    "segment_index": len(segments),
                    "source_kind": "asset",
                    "start_seconds": None,
                    "end_seconds": None,
                    "page_number": page_index,
                    "sheet_name": None,
                    "cell_range": None,
                    "slide_number": None,
                    "section_title": f"Page {page_index}",
                    "text_original": text,
                    "text_ptbr": None,
                    "language": metadata.get("language_original"),
                }
            )
    return segments


def docx_segments(metadata: dict[str, Any], file_path: Path) -> list[dict[str, Any]]:
    from docx import Document

    asset_id = metadata["asset_id"]
    document = Document(str(file_path))
    blocks: list[tuple[str | None, str]] = []
    current_title: str | None = None
    current_lines: list[str] = []

    def flush() -> None:
        nonlocal current_lines
        text = clean_text("\n".join(current_lines))
        if text:
            blocks.append((current_title, text))
        current_lines = []

    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            flush()
            continue
        style_name = (paragraph.style.name or "").lower()
        if "heading" in style_name or "titulo" in style_name:
            flush()
            current_title = text
            continue
        current_lines.append(text)
    flush()

    for table_index, table in enumerate(document.tables, start=1):
        rows = []
        for row in table.rows:
            cells = [clean_text(cell.text) for cell in row.cells]
            rows.append(" | ".join(cells))
        table_text = clean_text("\n".join(rows))
        if table_text:
            blocks.append((f"Table {table_index}", table_text))

    segments = []
    for index, (section_title, text) in enumerate(blocks):
        segments.append(
            {
                "segment_id": f"{asset_id}-asset-{index + 1:04d}",
                "segment_index": index,
                "source_kind": "asset",
                "start_seconds": None,
                "end_seconds": None,
                "page_number": None,
                "sheet_name": None,
                "cell_range": None,
                "slide_number": None,
                "section_title": section_title,
                "text_original": text,
                "text_ptbr": None,
                "language": metadata.get("language_original"),
            }
        )
    return segments


def xlsx_segments(metadata: dict[str, Any], file_path: Path) -> list[dict[str, Any]]:
    from openpyxl import load_workbook
    from openpyxl.utils import get_column_letter

    asset_id = metadata["asset_id"]
    workbook = load_workbook(str(file_path), data_only=True, read_only=True)
    segments = []
    for sheet in workbook.worksheets:
        rows = list(sheet.iter_rows(values_only=True))
        for row_index, row in enumerate(rows, start=1):
            values = ["" if value is None else str(value) for value in row]
            if not any(value.strip() for value in values):
                continue
            parts = [f"{get_column_letter(col_index)}: {value}" for col_index, value in enumerate(values, start=1) if value.strip()]
            text = clean_text("; ".join(parts))
            last_col = get_column_letter(max(len(values), 1))
            segments.append(
                {
                    "segment_id": f"{asset_id}-asset-{len(segments) + 1:04d}",
                    "segment_index": len(segments),
                    "source_kind": "asset",
                    "start_seconds": None,
                    "end_seconds": None,
                    "page_number": None,
                    "sheet_name": sheet.title,
                    "cell_range": f"A{row_index}:{last_col}{row_index}",
                    "slide_number": None,
                    "section_title": f"{sheet.title} row {row_index}",
                    "text_original": text,
                    "text_ptbr": None,
                    "language": metadata.get("language_original"),
                }
            )
    workbook.close()
    return segments


def pptx_segments(metadata: dict[str, Any], file_path: Path) -> list[dict[str, Any]]:
    from pptx import Presentation

    asset_id = metadata["asset_id"]
    presentation = Presentation(str(file_path))
    segments = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = clean_text(shape.text)
                if text:
                    texts.append(text)
        slide_text = clean_text("\n".join(texts))
        if not slide_text:
            continue
        segments.append(
            {
                "segment_id": f"{asset_id}-asset-{len(segments) + 1:04d}",
                "segment_index": len(segments),
                "source_kind": "asset",
                "start_seconds": None,
                "end_seconds": None,
                "page_number": None,
                "sheet_name": None,
                "cell_range": None,
                "slide_number": slide_index,
                "section_title": f"Slide {slide_index}",
                "text_original": slide_text,
                "text_ptbr": None,
                "language": metadata.get("language_original"),
            }
        )
    return segments


def find_original_file(metadata: dict[str, Any], metadata_path: Path) -> Path:
    storage_path = Path(metadata["storage_path"])
    if storage_path.exists():
        return storage_path
    candidate = metadata_path.parent / storage_path.name
    if candidate.exists():
        return candidate
    raise FileNotFoundError(f"Could not find original file for asset {metadata['asset_id']}: {storage_path}")


def process_asset(metadata_path: Path, output_path: Path) -> dict[str, Any]:
    metadata = load_json(metadata_path)
    file_path = find_original_file(metadata, metadata_path)
    asset_type = metadata.get("asset_type")

    suffix = file_path.suffix.lower()
    if asset_type == "pdf" and suffix == ".pdf":
        segments = pdf_segments(metadata, file_path)
    elif asset_type == "doc" and suffix == ".docx":
        segments = docx_segments(metadata, file_path)
    elif asset_type == "spreadsheet" and suffix == ".csv":
        segments = csv_segments(metadata, file_path)
    elif asset_type == "spreadsheet" and suffix in {".xlsx", ".xlsm"}:
        segments = xlsx_segments(metadata, file_path)
    elif asset_type == "slides" and suffix == ".pptx":
        segments = pptx_segments(metadata, file_path)
    elif asset_type in {"text", "html"} or suffix in {".txt", ".md", ".html", ".htm"}:
        segments = text_segments(metadata, file_path)
    else:
        raise ValueError(
            f"Asset type {asset_type!r} with suffix {file_path.suffix!r} is not supported yet by process_asset.py"
        )

    payload = {
        "schema_version": "1.0",
        "episode_video_id": metadata["episode_video_id"],
        "asset_id": metadata["asset_id"],
        "source_kind": "asset",
        "language_original": metadata.get("language_original"),
        "segments": segments,
    }
    write_json(output_path, payload)
    return payload


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--metadata", required=True, type=Path, help="Path to data/raw/assets/{asset_id}/metadata.json")
    parser.add_argument("--output", type=Path, help="Path to write content_segments.json")
    args = parser.parse_args()

    metadata = load_json(args.metadata)
    output_path = args.output or Path("data/processed/assets") / metadata["asset_id"] / "content_segments.json"
    payload = process_asset(args.metadata, output_path)
    print(f"Wrote {len(payload['segments'])} asset segment(s) to {output_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
